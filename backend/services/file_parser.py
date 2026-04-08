"""
services/file_parser.py — Centralized file-to-text extraction
From chart.svg: FileParser → DeepDoc解析引擎 → multiple format parsers

Supported formats:
  PDF    — pdfplumber
  DOCX   — python-docx (paragraphs + tables)
  TXT/MD — encoding-aware plain text
  XLSX   — openpyxl (all sheets)
  XLS    — openpyxl (falls back gracefully)
  HTML   — stdlib HTMLParser (strips tags/scripts)
  JSON   — recursive string extraction
  PPTX   — python-pptx (slide text frames)
"""

import io
import json
from html.parser import HTMLParser


# Extensions this module can handle
SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "doc",
    "txt", "md",
    "xlsx", "xls",
    "html", "htm",
    "json",
    "pptx",
}


def parse_file(filename: str, data: bytes) -> str:
    """Dispatch to the correct parser based on file extension."""
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    _parsers = {
        "pdf":  _parse_pdf,
        "docx": _parse_docx,
        "doc":  _parse_docx,   # python-docx handles many .doc files
        "txt":  _parse_text,
        "md":   _parse_text,
        "xlsx": _parse_excel,
        "xls":  _parse_excel,
        "html": _parse_html,
        "htm":  _parse_html,
        "json": _parse_json,
        "pptx": _parse_pptx,
    }
    return _parsers.get(ext, _parse_text)(data)


# ── Individual parsers ────────────────────────────────────────────────────────

def _parse_pdf(data: bytes) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _parse_docx(data: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(data))
    parts = []
    # Paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    # Tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(
                cell.text.strip() for cell in row.cells if cell.text.strip()
            )
            if row_text:
                parts.append(row_text)
    return "\n".join(parts)


def _parse_text(data: bytes) -> str:
    """Try common encodings; latin-1 is the final fallback (never raises)."""
    for encoding in ("utf-8-sig", "utf-8", "utf-16", "gb18030", "latin-1"):
        try:
            return data.decode(encoding)
        except (UnicodeDecodeError, Exception):
            continue
    return data.decode("latin-1")


def _parse_excel(data: bytes) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for sheet in wb.worksheets:
        parts.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            row_text = "\t".join(
                str(cell) if cell is not None else "" for cell in row
            )
            if row_text.strip():
                parts.append(row_text)
    wb.close()
    return "\n".join(parts)


class _HTMLTextExtractor(HTMLParser):
    """Strip tags and skip script/style blocks."""

    _SKIP = {"script", "style", "head", "meta", "link", "noscript"}

    def __init__(self):
        super().__init__()
        self._parts: list[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag.lower() in self._SKIP:
            self._skip = True

    def handle_endtag(self, tag):
        if tag.lower() in self._SKIP:
            self._skip = False

    def handle_data(self, data):
        if not self._skip:
            stripped = data.strip()
            if stripped:
                self._parts.append(stripped)

    def get_text(self) -> str:
        return "\n".join(self._parts)


def _parse_html(data: bytes) -> str:
    raw = _parse_text(data)   # encoding detection first
    extractor = _HTMLTextExtractor()
    extractor.feed(raw)
    return extractor.get_text()


def _collect_strings(obj, parts: list[str]) -> None:
    if isinstance(obj, str):
        if obj.strip():
            parts.append(obj.strip())
    elif isinstance(obj, dict):
        for v in obj.values():
            _collect_strings(v, parts)
    elif isinstance(obj, list):
        for item in obj:
            _collect_strings(item, parts)


def _parse_json(data: bytes) -> str:
    raw = _parse_text(data)
    try:
        obj = json.loads(raw)
        parts: list[str] = []
        _collect_strings(obj, parts)
        return "\n".join(parts) if parts else raw
    except json.JSONDecodeError:
        return raw   # not valid JSON — return raw text


def _parse_pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts.append(f"[Slide {slide_num}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)
